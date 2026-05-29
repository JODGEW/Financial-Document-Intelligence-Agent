"""PowerPoint format handler.

Uses python-pptx directly so we can map slides to Documents one-to-one and
append table content in the same ``[Extracted Tables]`` format other handlers
use. Avoids ``UnstructuredPowerPointLoader`` to keep the dependency surface
small (no nltk_data download on first ingest).
"""

from __future__ import annotations

from pathlib import Path

from langchain_core.documents import Document

from .registry import FormatHandler, register
from .tables import extract_pptx_tables


def _slide_title(slide) -> str | None:
    """Return the slide's title text, or None if no title placeholder is set."""
    try:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip()
            return title or None
    except Exception:
        return None
    return None


def _slide_body_text(slide) -> str:
    """Collect text from non-title shapes on a slide."""
    title_shape_id = None
    try:
        if slide.shapes.title is not None:
            title_shape_id = slide.shapes.title.shape_id
    except Exception:
        title_shape_id = None

    lines: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "shape_id", None) == title_shape_id:
            continue
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if text:
                    lines.append(text)
    return "\n".join(lines)


def _load(path: str) -> list[Document]:
    from pptx import Presentation

    presentation = Presentation(path)
    tables_by_slide = extract_pptx_tables(path)

    docs: list[Document] = []
    for slide_index, slide in enumerate(presentation.slides):
        title = _slide_title(slide)
        body = _slide_body_text(slide)
        if not body and not title and slide_index not in tables_by_slide:
            continue

        slide_number = slide_index + 1
        section_title = (
            f"Slide {slide_number}: {title}" if title else f"Slide {slide_number}"
        )
        parts: list[str] = []
        if title:
            parts.append(title)
        if body:
            parts.append(body)
        page_content = "\n\n".join(parts).strip() or section_title

        metadata: dict = {
            "source": path,
            "slide_number": slide_number,
            "section_title": section_title,
        }

        table_text = tables_by_slide.get(slide_index)
        if table_text:
            page_content = (
                f"{page_content.rstrip()}\n\n[Extracted Tables]\n{table_text}"
            )
            metadata["contains_tables"] = True

        docs.append(Document(page_content=page_content, metadata=metadata))

    if not docs and tables_by_slide:
        # Edge case: deck has only tables and no text. Emit one Document per
        # slide so retrieval can still surface table content.
        for slide_index, table_text in sorted(tables_by_slide.items()):
            slide_number = slide_index + 1
            section_title = f"Slide {slide_number}"
            docs.append(
                Document(
                    page_content=f"{section_title}\n\n[Extracted Tables]\n{table_text}",
                    metadata={
                        "source": path,
                        "slide_number": slide_number,
                        "section_title": section_title,
                        "contains_tables": True,
                    },
                )
            )
    return docs


def _split(docs: list[Document]) -> list[Document]:
    from ingest import _recursive_splitter

    splitter = _recursive_splitter(chunk_size=900, chunk_overlap=100)
    chunks: list[Document] = []
    for doc in docs:
        for piece in splitter.split_documents([doc]):
            piece.metadata = {**doc.metadata, **piece.metadata}
            chunks.append(piece)
    return chunks


register(
    FormatHandler(
        extensions=(".pptx",),
        loader=_load,
        splitter=_split,
        format_family="text",
        extract_tables=True,
    )
)
