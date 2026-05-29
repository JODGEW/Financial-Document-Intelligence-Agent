"""Generic table-formatting helpers shared across format handlers.

Per-format table extraction (PDF, DOCX, HTML, XLSX, PPTX) ultimately funnels
into ``format_table_rows`` so retrieval sees a consistent markdown-like
representation.
"""

from __future__ import annotations

from pathlib import Path


def format_table_rows(rows: list[list[object]]) -> str:
    """Format extracted table rows as markdown-like text for retrieval."""
    cleaned_rows = [
        ["" if cell is None else " ".join(str(cell).split()) for cell in row]
        for row in rows
        if any(cell is not None and str(cell).strip() for cell in row)
    ]
    if not cleaned_rows:
        return ""

    width = max(len(row) for row in cleaned_rows)
    normalized = [row + [""] * (width - len(row)) for row in cleaned_rows]
    lines = ["| " + " | ".join(row) + " |" for row in normalized]
    if len(lines) > 1:
        lines.insert(1, "| " + " | ".join(["---"] * width) + " |")
    return "\n".join(lines)


def extract_docx_tables(source_path: str | Path) -> list[str]:
    """Extract every <w:tbl> in a .docx as markdown-like text blocks."""
    from docx import Document as DocxDocument  # python-docx

    formatted_tables: list[str] = []
    try:
        document = DocxDocument(str(source_path))
    except Exception as exc:
        print(f"  WARN: failed to open DOCX {Path(source_path).name}: {exc}")
        return formatted_tables

    for index, table in enumerate(document.tables, 1):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        formatted = format_table_rows(rows)
        if formatted:
            formatted_tables.append(f"Table {index}\n{formatted}")
    return formatted_tables


def extract_html_tables(html_text: str) -> list[str]:
    """Extract every <table> in an HTML string as markdown-like text blocks."""
    from bs4 import BeautifulSoup

    formatted_tables: list[str] = []
    soup = BeautifulSoup(html_text, "lxml")
    for index, table in enumerate(soup.find_all("table"), 1):
        rows: list[list[str]] = []
        for tr in table.find_all("tr"):
            cells = tr.find_all(["th", "td"])
            rows.append([cell.get_text(" ", strip=True) for cell in cells])
        formatted = format_table_rows(rows)
        if formatted:
            formatted_tables.append(f"Table {index}\n{formatted}")
    return formatted_tables


def extract_pptx_tables(source_path: str | Path) -> dict[int, str]:
    """Extract tables per slide from a .pptx file.

    Returns ``{slide_index: joined_markdown_tables}`` where ``slide_index`` is
    zero-based. Slides without tables are absent from the dict.
    """
    from pptx import Presentation

    tables_by_slide: dict[int, list[str]] = {}
    try:
        presentation = Presentation(str(source_path))
    except Exception as exc:
        print(f"  WARN: failed to open PPTX {Path(source_path).name}: {exc}")
        return {}

    for slide_index, slide in enumerate(presentation.slides):
        slide_blocks: list[str] = []
        table_index = 0
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            rows = [
                [cell.text for cell in row.cells]
                for row in table.rows
            ]
            formatted = format_table_rows(rows)
            if formatted:
                table_index += 1
                slide_blocks.append(f"Table {table_index}\n{formatted}")
        if slide_blocks:
            tables_by_slide[slide_index] = "\n\n".join(slide_blocks)
    return tables_by_slide
